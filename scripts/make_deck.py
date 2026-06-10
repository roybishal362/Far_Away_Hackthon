"""Generate the Kakehashi submission deck (.pptx) with real diagrams, a chart, and speaker notes.
Run: python scripts/make_deck.py  ->  Kakehashi_FAR_AWAY_2026.pptx
"""
import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "Kakehashi_FAR_AWAY_2026.pptx")

INDIGO = RGBColor(0x1F, 0x21, 0x50)
INDIGO2 = RGBColor(0x2B, 0x2D, 0x6E)
SAKURA = RGBColor(0xE1, 0x1D, 0x54)
SAKTINT = RGBColor(0xFC, 0xE7, 0xED)
MARI = RGBColor(0xFF, 0x95, 0x00)
MARTINT = RGBColor(0xFF, 0xF1, 0xDE)
EMER = RGBColor(0x10, 0xB9, 0x81)
EMTINT = RGBColor(0xE3, 0xF7, 0xF0)
INK = RGBColor(0x15, 0x16, 0x3A)
GREY = RGBColor(0x70, 0x70, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INDTINT = RGBColor(0xEC, 0xED, 0xF5)
PETAL = RGBColor(0xFF, 0x9B, 0xB3)
LINE = RGBColor(0xD8, 0xD8, 0xE2)

C, M, T, B = PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, MSO_ANCHOR.TOP, PP_ALIGN.LEFT

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def _set(shape, text, size=14, color=INK, bold=False, align=C, anchor=M):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Pt(7))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, Pt(3))
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        t, o = ln if isinstance(ln, tuple) else (ln, {})
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(o.get("size", size)); f.bold = o.get("bold", bold)
        f.color.rgb = o.get("color", color); f.name = "Calibri"


def rect(s, l, t, w, h, fill, rounded=True, line=None, text=None, **kw):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if text is not None:
        _set(shp, text, **kw)
    return shp


def shape(s, kind, l, t, w, h, fill, text=None, **kw):
    shp = s.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    if text is not None:
        _set(shp, text, **kw)
    return shp


def txt(s, l, t, w, h, text, **kw):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _set(tb, text, **kw)
    return tb


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def base(title, accent=SAKURA, n=None, kicker=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, -0.1, -0.1, 13.6, 0.2, accent, rounded=False)
    if kicker:
        txt(s, 0.6, 0.42, 11, 0.35, kicker, size=12, color=accent, bold=True, align=B)
    txt(s, 0.6, 0.72 if kicker else 0.55, 12.1, 1.0, title, size=30, color=INDIGO, bold=True, align=B)
    txt(s, 0.6, 7.02, 9, 0.4, "Kakehashi 架け橋  ·  FAR AWAY 2026  ·  Agentic & Autonomous Systems", size=10, color=GREY, align=B)
    if n:
        txt(s, 12.4, 7.02, 0.6, 0.4, str(n), size=10, color=GREY, align=PP_ALIGN.RIGHT)
    return s


# ───────────────────────── 1 · TITLE ─────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, -0.1, -0.1, 13.6, 7.7, INDIGO, rounded=False)
# bridge motif: two pillars + an arc
rect(s, 1.4, 3.2, 0.5, 2.4, PETAL)
rect(s, 11.5, 3.2, 0.5, 2.4, PETAL)
arc = s.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(1.4), Inches(2.0), Inches(10.6), Inches(2.6), )
arc.fill.solid(); arc.fill.fore_color.rgb = SAKURA; arc.line.fill.background(); arc.shadow.inherit = False
txt(s, 1.2, 0.6, 4, 0.5, "India  🇮🇳", size=16, color=PETAL, bold=True, align=B)
txt(s, 9.0, 0.6, 4, 0.5, "🇯🇵  Japan", size=16, color=PETAL, bold=True, align=PP_ALIGN.RIGHT)
txt(s, 0.8, 4.55, 11.7, 1.2, [("Kakehashi ", {"size": 60, "bold": True, "color": WHITE}), ("架け橋", {"size": 42, "color": PETAL})], align=C)
txt(s, 0.8, 5.75, 11.7, 0.6, "An autonomous AI bridge for Indian workers to Japan", size=23, color=WHITE, align=C)
txt(s, 0.8, 6.5, 11.7, 0.5, "Real data · Cited sources · Measured proof · EN / हिन्दी / 日本語   |   Theme: Agentic & Autonomous Systems",
    size=13, color=RGBColor(0xCF, 0xD0, 0xE0), align=C)
notes(s, "One line: Kakehashi is an autonomous AI that guides Indian workers through Japan's Specified Skilled Worker visa journey. "
         "Hook the dual-government India-Japan corridor and that everything is real and cited. Keep it to 8 seconds.")

# ───────────────────────── 2 · PROBLEM ─────────────────────────
s = base("The problem", SAKURA, 2, kicker="MOTIVATION")
txt(s, 0.6, 1.55, 12.1, 0.5, "A worker who wants to go to Japan faces a scattered, months-long maze:", size=16, color=INK)
pains = [("❓ Which visa?", "Am I even eligible?"), ("🗣️ JLPT N4 bar", "+ sector skills tests"),
         ("📑 Paperwork", "CoE, embassy, employer"), ("💸 Middlemen", "scams & misinformation")]
for i, (a, b) in enumerate(pains):
    x = 0.6 + i * 3.1
    rect(s, x, 2.3, 2.85, 1.7, SAKTINT, text=[(a, {"size": 16, "bold": True, "color": SAKURA}), (b, {"size": 12, "color": INK})])
rect(s, 0.6, 4.5, 12.1, 1.7, INDTINT,
     text=[("The information is real and public — but buried across Japanese government sites.", {"size": 16, "bold": True, "color": INDIGO}),
           ("No single, trustworthy, personalized guide exists. That gap is what we close.", {"size": 14, "color": INK})])
notes(s, "Make it human: imagine a nurse in Delhi who qualifies but has no idea where to start. The info exists but is fragmented and gatekept. "
         "This is a real pain for real people — set up the impact.")

# ───────────────────────── 3 · OPPORTUNITY ─────────────────────────
s = base("Why now — a dual-government opportunity", INDIGO, 3, kicker="REAL-WORLD IMPACT")
stats = [("50,000", "skilled Indian workers → Japan\n(signed Aug-2025 pact)", SAKURA),
         ("570,000", "care-worker shortfall Japan faces\nby 2040 (MHLW)", MARI),
         ("AI", "named by BOTH govts as the tool\nfor 'trustworthy talent mobility'", EMER)]
for i, (big, sub, col) in enumerate(stats):
    x = 0.7 + i * 4.05
    rect(s, x, 2.2, 3.75, 2.7, INDTINT)
    txt(s, x, 2.5, 3.75, 1.1, big, size=46, color=col, bold=True, align=C)
    txt(s, x, 3.7, 3.75, 1.1, sub.replace("\n", "  "), size=13, color=INK, align=C)
rect(s, 0.7, 5.2, 11.95, 1.1, INDIGO,
     text=[("Kakehashi serves the exact corridor India & Japan just committed to — and the finale is in Japan.", {"size": 15, "bold": True, "color": WHITE})])
notes(s, "These are real, current, sourced numbers. The strategic point: we're not inventing a problem — we're serving a corridor two "
         "governments signed in 2025, judged at a finale in Japan. That's deliberate.")

# ───────────────────────── 4 · SOLUTION FLOW ─────────────────────────
s = base("What Kakehashi does", EMER, 4, kicker="THE SOLUTION")
txt(s, 0.6, 1.5, 12, 0.5, "You tell it about you (or upload your resume). Autonomous agents do the rest:", size=15, color=INK)
flow = [("👤", "You /\nresume"), ("🧭", "Visa\npathway"), ("💼", "Real\njobs"), ("📄", "Procedure\n+ links"), ("📚", "Study\nplan"), ("💴", "Salary\n& cost")]
for i, (ic, lab) in enumerate(flow):
    x = 0.5 + i * 2.05
    rect(s, x, 2.5, 1.7, 1.5, INDTINT if i else SAKTINT,
         text=[(ic, {"size": 22}), (lab.replace("\n", " "), {"size": 11, "color": INK, "bold": True})])
    if i < len(flow) - 1:
        shape(s, MSO_SHAPE.RIGHT_ARROW, x + 1.72, 3.05, 0.32, 0.4, MARI)
rect(s, 0.5, 4.6, 12.3, 1.6, EMTINT,
     text=[("…then ask follow-ups in a grounded chat — in your language — and save / share / download the whole plan as a PDF dossier.",
            {"size": 15, "color": INK, "bold": True})])
notes(s, "Walk left to right. Emphasize it's one flow, fully automated after the form. Mention the chat + PDF dossier as the 'take it with you' moment.")

# ───────────────────────── 5 · ARCHITECTURE ─────────────────────────
s = base("Architecture", INDIGO, 5, kicker="ENGINEERING DEPTH")
rect(s, 0.7, 1.7, 11.95, 0.9, INDTINT, text=[("FRONTEND — Next.js (Home · How it works · Build my plan)", {"size": 14, "bold": True, "color": INDIGO})])
shape(s, MSO_SHAPE.DOWN_ARROW, 6.4, 2.65, 0.5, 0.45, GREY)
txt(s, 7.0, 2.62, 5, 0.5, "live agent timeline over SSE", size=11, color=GREY, align=B)
rect(s, 0.7, 3.2, 11.95, 0.9, INDIGO, text=[("FastAPI  —  /run · /run/stream (SSE) · /eval · /chat · /resume · /dossier · /save · /health",
                                             {"size": 13, "bold": True, "color": WHITE})])
shape(s, MSO_SHAPE.DOWN_ARROW, 6.4, 4.15, 0.5, 0.45, GREY)
rect(s, 0.7, 4.7, 11.95, 1.85, RGBColor(0xF4, 0xF1, 0xF7), line=LINE)
txt(s, 0.9, 4.78, 6, 0.4, "core/ — UI-agnostic Python brain (engine orchestrates + adapts)", size=12, color=INDIGO, bold=True, align=B)
core = [("🤖 6 Agents", SAKTINT, SAKURA), ("🔌 Real tools\nJSearch · e-Stat", INDTINT, INDIGO),
        ("📚 RAG\ncited facts", EMTINT, EMER), ("🧪 Eval\nablation proof", MARTINT, MARI)]
for i, (lab, fill, col) in enumerate(core):
    x = 0.95 + i * 2.92
    rect(s, x, 5.25, 2.7, 1.05, fill, text=[(lab.replace("\n", " "), {"size": 12, "bold": True, "color": col})])
notes(s, "Three clean layers. The point for technical judges: the brain is UI-agnostic and tool-based — every agent calls a REAL tool. "
         "The SSE stream is why you can watch agents think live.")

# ───────────────────────── 6 · AGENTS ─────────────────────────
s = base("The multi-agent system", SAKURA, 6, kicker="HOW IT THINKS")
txt(s, 0.6, 1.5, 12, 0.45, "Each agent:  reason  →  call a real tool  →  cite the source  →  score confidence", size=15, color=INDIGO, bold=True)
agents = [("🧭 Pathway", "right visa + verdict, readiness, gaps (cited)"),
          ("💼 Jobs", "real live openings, ranked by fit to you"),
          ("📄 Procedure", "official steps, each with a real gov link"),
          ("📚 Prep", "study plan for your level + free resources"),
          ("💴 Synthesis", "your overview + salary + cost"),
          ("💬 Chat", "grounded follow-up Q&A, with citations")]
for i, (a, b) in enumerate(agents):
    x = 0.6 + (i % 3) * 4.05
    y = 2.25 + (i // 3) * 1.9
    rect(s, x, y, 3.8, 1.65, INDTINT, line=LINE,
         text=[(a, {"size": 16, "bold": True, "color": SAKURA}), (b, {"size": 12, "color": INK})])
notes(s, "Don't read all six — say 'six specialized agents, each grounded and citing sources' and point at Pathway + Jobs as examples. "
         "Contrast with a single-prompt 'AI wrapper', which FAR AWAY explicitly penalizes.")

# ───────────────────────── 7 · AUTONOMY ─────────────────────────
s = base("Autonomy in action — the plan adapts to YOU", MARI, 7, kicker="THE 'WOW'")
rect(s, 5.1, 1.7, 3.1, 0.9, INDIGO, text=[("Profile → Orchestrator", {"size": 14, "bold": True, "color": WHITE}), ("classifies your visa route", {"size": 10, "color": PETAL})])
branches = [("👩‍⚕️ Nurse", "SSW caregiving — FULL plan:\njobs · procedure · study · cost", EMER, EMTINT),
            ("👨‍💻 Software eng.", "Reroutes to ENGINEER visa —\nSSW-only steps auto-skipped", SAKURA, SAKTINT),
            ("🧑‍💼 HR / office", "Reroutes to SPECIALIST visa\n(技人国)", MARI, MARTINT)]
for i, (a, b, col, fill) in enumerate(branches):
    x = 0.7 + i * 4.05
    shape(s, MSO_SHAPE.DOWN_ARROW, x + 1.65, 2.7, 0.4, 0.45, GREY)
    rect(s, x, 3.25, 3.75, 2.0, fill, line=col,
         text=[(a, {"size": 16, "bold": True, "color": col}), (b.replace("\n", " "), {"size": 12, "color": INK})])
rect(s, 0.7, 5.5, 11.95, 0.95, INDIGO,
     text=[("Different person → different plan. The agent DECIDES — and even refuses SSW when it doesn't fit. That's real autonomy.", {"size": 14, "bold": True, "color": WHITE})])
notes(s, "This is your strongest live moment. In the demo, click the Nurse persona then the Software persona and let judges SEE the plan "
         "change shape and skip steps. Most projects force everyone down one path; ours refuses and reroutes.")

# ───────────────────────── 8 · PROOF (chart) ─────────────────────────
s = base("Proof, not claims", EMER, 8, kicker="MEASURED")
txt(s, 0.6, 1.5, 12, 0.45, "Same question, grounded (our agents) vs. ungrounded (plain LLM), scored against a gold set of official facts.", size=14, color=INK)
cd = CategoryChartData()
cd.categories = ["Factual accuracy (%)"]
cd.add_series("Grounded (ours)", (86,))
cd.add_series("Ungrounded LLM (avg)", (48,))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.2), Inches(6.6), Inches(4.4), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = EMER
ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb = SAKURA
rect(s, 7.7, 2.4, 4.9, 1.9, EMTINT, text=[("0", {"size": 44, "bold": True, "color": EMER}), ("hallucinations — grounded, EVERY run", {"size": 13, "color": INK})])
rect(s, 7.7, 4.5, 4.9, 1.9, SAKTINT, text=[("29-71%", {"size": 34, "bold": True, "color": SAKURA}), ("plain-LLM accuracy swings run-to-run (ours holds a steady 86%)", {"size": 13, "color": INK})])
notes(s, "Honest & reproducible (PROOF.md, 3 runs): grounded = 86% accuracy + 0 hallucinations EVERY run; the plain LLM swings 29-71%. "
         "The win is accuracy AND consistency. Measured live in the Proof tab.")

# ───────────────────────── 9 · REAL DATA ─────────────────────────
s = base("Built on real data", INDIGO, 9, kicker="REAL OR NOTHING")
rows = [("SSW rules & procedures", "official ssw.go.jp (ISA) + MOFA", "grounded & cited (RAG)"),
        ("Live job openings", "JSearch real-time API", "ranked by fit"),
        ("Government labour stats", "Japan e-Stat official API", "authoritative"),
        ("Fees & salaries", "Prometric · JLPT · JFT + research", "sourced")]
for i, (a, b, c) in enumerate(rows):
    y = 2.0 + i * 1.05
    rect(s, 0.7, y, 4.0, 0.9, INDTINT, text=[(a, {"size": 13, "bold": True, "color": INDIGO})])
    shape(s, MSO_SHAPE.RIGHT_ARROW, 4.8, y + 0.27, 0.35, 0.36, MARI)
    rect(s, 5.3, y, 4.6, 0.9, RGBColor(0xF4, 0xF1, 0xF7), line=LINE, text=[(b, {"size": 12, "color": INK})])
    rect(s, 10.0, y, 2.6, 0.9, EMTINT, text=[(c, {"size": 11, "bold": True, "color": EMER})])
txt(s, 0.7, 6.35, 12, 0.5, "No key → the tool says 'not configured'. It never invents data.", size=14, color=SAKURA, bold=True)
notes(s, "Hammer 'real or nothing'. Every claim traces to a real source; missing data is shown honestly, never faked — which is exactly "
         "what FAR AWAY rewards and what most demos fail.")

# ───────────────────────── 10 · PRODUCT ─────────────────────────
s = base("A real product, not a demo", SAKURA, 10, kicker="DESIGN & EXECUTION")
txt(s, 0.6, 1.5, 12, 0.4, "Multi-page site · 8 result tabs · resume auto-fill · save/share link · PDF dossier · progress checklist", size=14, color=INK)
tabs = ["Overview", "Pathway", "Jobs", "Procedure", "Study plan", "Journey", "Proof", "Ask AI"]
for i, t in enumerate(tabs):
    x = 0.7 + (i % 4) * 3.05
    y = 2.2 + (i // 4) * 0.95
    rect(s, x, y, 2.85, 0.75, INDTINT, line=LINE, text=[(t, {"size": 13, "bold": True, "color": INDIGO})])
rect(s, 0.7, 4.35, 5.85, 1.9, EMTINT, text=[("🌐 Multilingual", {"size": 18, "bold": True, "color": EMER}), ("Content AND chat in English / हिन्दी / 日本語 — for the Japan finale", {"size": 13, "color": INK})])
rect(s, 6.75, 4.35, 5.85, 1.9, MARTINT, text=[("🔐 Private by design", {"size": 18, "bold": True, "color": MARI}), ("Encryption at rest · no PII in logs · honest disclaimer", {"size": 13, "color": INK})])
notes(s, "Show it looks shipped, not hacked. The 日本語 toggle is a deliberate gift to the Japanese finale judges. Mention privacy as a trust signal.")

# ───────────────────────── 11 · IMPACT (bridge) ─────────────────────────
s = base("A true bridge — who it helps", INDIGO, 11, kicker="REAL-WORLD IMPACT")
rect(s, 0.7, 2.2, 3.3, 2.6, SAKTINT, text=[("🇮🇳 India", {"size": 20, "bold": True, "color": SAKURA}), ("Formal overseas jobs + remittances; an alternative to scam middlemen", {"size": 13, "color": INK})])
rect(s, 9.3, 2.2, 3.3, 2.6, INDTINT, text=[("🇯🇵 Japan", {"size": 20, "bold": True, "color": INDIGO}), ("Fills its designated-sector labour shortage with vetted talent", {"size": 13, "color": INK})])
arc2 = s.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Inches(4.0), Inches(2.0), Inches(5.3), Inches(1.7))
arc2.fill.solid(); arc2.fill.fore_color.rgb = MARI; arc2.line.fill.background(); arc2.shadow.inherit = False
rect(s, 4.4, 3.55, 4.5, 1.0, INDIGO, text=[("Kakehashi", {"size": 18, "bold": True, "color": WHITE}), ("the cited, honest path", {"size": 11, "color": PETAL})])
rect(s, 0.7, 5.2, 11.95, 1.1, EMTINT, text=[("The worker gets a trustworthy plan + protection from recruitment-fee scams (official FREE job-matching).", {"size": 14, "bold": True, "color": INK})])
notes(s, "Tie it back to the dual-government pact. The anti-exploitation angle (free official job-matching, sending orgs optional for India) "
         "is a genuine social-impact story — lead with the worker.")

# ───────────────────────── 12 · TECH & ROADMAP ─────────────────────────
s = base("Engineering, scalability & honest roadmap", MARI, 12, kicker="BUILT TO LAST")
txt(s, 0.6, 1.55, 12, 0.4, "Next.js · FastAPI (SSE) · Groq gpt-oss-120b (multilingual) · BM25 RAG · Fernet (AES)", size=14, color=INDIGO, bold=True)
rect(s, 0.7, 2.3, 11.95, 1.0, INDTINT, text=[("UI-agnostic core with swappable agents & tools — add a sector, agent, or visa route without a rewrite.", {"size": 14, "color": INK})])
rect(s, 0.7, 3.6, 5.85, 2.5, EMTINT)
txt(s, 0.95, 3.8, 5.4, 0.4, "✅ Live now", size=16, bold=True, color=EMER, align=B)
txt(s, 0.95, 4.3, 5.4, 1.7, [("Pathway · Jobs (real) · Procedure · Prep · Synthesis", {"size": 12}), ("Chat · Proof · save/share · PDF · EN/HI/JA", {"size": 12}), ("Adaptive visa routing", {"size": 12})], color=INK, align=B)
rect(s, 6.75, 3.6, 5.85, 2.5, MARTINT)
txt(s, 7.0, 3.8, 5.4, 0.4, "🔜 Roadmap (wired, marked)", size=16, bold=True, color=MARI, align=B)
txt(s, 7.0, 4.3, 5.4, 1.7, [("Amadeus flights (paid tier)", {"size": 12}), ("e-Stat 'Demand' charts", {"size": 12}), ("Recruiter-scam detection", {"size": 12})], color=INK, align=B)
notes(s, "Two messages: (1) it's engineered to extend, not a throwaway; (2) we're honest about what's live vs roadmap — judges trust honest "
         "roadmapping and catch the opposite.")

# ───────────────────────── 13 · CLOSING ─────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, -0.1, -0.1, 13.6, 7.7, INDIGO, rounded=False)
rect(s, 0, 6.7, 13.4, 0.9, MARI, rounded=False)
txt(s, 0.85, 2.2, 11.6, 1.1, "A bridge between India and Japan.", size=40, color=WHITE, bold=True)
txt(s, 0.9, 3.5, 11.5, 1.8,
    [("Real data. Cited sources. Measured proof. Genuinely autonomous.", {"size": 22, "color": PETAL}),
     ("🔗  github.com/roybishal362/Far_Away_Hackthon", {"size": 18, "color": WHITE}),
     ("Working demo  ·  add your live URL here", {"size": 14, "color": RGBColor(0xCF, 0xD0, 0xE0)})], align=B)
txt(s, 0.85, 6.78, 11.6, 0.6, "Kakehashi 架け橋  —  FAR AWAY 2026 · Agentic & Autonomous Systems", size=14, color=WHITE, bold=True, align=B)
notes(s, "Close on the bridge metaphor and the four words. Invite them to try the live demo / scan the GitHub. End in under 8 seconds.")

try:
    prs.save(OUT)
    print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
except PermissionError:
    alt = OUT.replace(".pptx", "_NEW.pptx")
    prs.save(alt)
    print("Original was open/locked. Saved upgraded deck to:", alt, "| slides:", len(prs.slides._sldIdLst))
